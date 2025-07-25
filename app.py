import streamlit as st
import boto3
import PyPDF2
import io
import json
from datetime import datetime
from tavily import TavilyClient
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def init_bedrock_client():
    """AWS Bedrockクライアントを初期化"""
    try:
        bedrock_client = boto3.client(
            'bedrock-runtime',
            region_name=st.secrets["aws"]["AWS_REGION"],
            aws_access_key_id=st.secrets["aws"]["AWS_ACCESS_KEY_ID"],
            aws_secret_access_key=st.secrets["aws"]["AWS_SECRET_ACCESS_KEY"]
        )
        return bedrock_client
    except Exception as e:
        st.error(f"AWS Bedrock接続エラー: {e}")
        return None

def extract_text_from_pdf(pdf_file):
    """PDFファイルからテキストを抽出"""
    try:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        st.error(f"PDF読み込みエラー: {e}")
        return None

def extract_text_from_pptx(pptx_file):
    """PowerPointファイルからテキストを抽出"""
    try:
        # ファイルポインタを先頭に戻す
        pptx_file.seek(0)
        
        # PowerPointファイルを読み込み
        presentation = Presentation(pptx_file)
        text = ""
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            text += f"\n--- スライド {slide_num} ---\n"
            
            # スライド内の全ての図形からテキストを抽出
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
                
                # 表がある場合のテキスト抽出
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text:
                                row_text.append(cell.text.strip())
                        if row_text:
                            text += " | ".join(row_text) + "\n"
            
            text += "\n"
        
        return text
    except Exception as e:
        st.error(f"PowerPoint読み込みエラー: {e}")
        return None

def init_tavily_client():
    """Tavily APIクライアントを初期化"""
    try:
        tavily_client = TavilyClient(api_key=st.secrets["tavily"]["API_KEY"])
        return tavily_client
    except Exception as e:
        st.error(f"Tavily API接続エラー: {e}")
        return None

# システムプロンプト定義
KEYWORD_EXTRACTION_PROMPT_TEMPLATE = """
以下の決裁書から、関連情報を検索するための効果的なキーワードを抽出してください。

【決裁書内容】
{document_text}

【指示】
1. この決裁書の内容に最も関連する検索キーワードを3個抽出
2. 各キーワードは具体的で検索に適したものにする
3. 業界情報、法規制、ベストプラクティスを見つけるのに有効なキーワード
4. 結果は以下の形式で出力：

キーワード1: [具体的なキーワード]
キーワード2: [具体的なキーワード]
キーワード3: [具体的なキーワード]

※キーワードのみを出力し、説明は不要です。
"""

DEFAULT_REVIEW_PROMPT_TEMPLATE = """あなたは製造業の情シス部門の経験豊富な上司として振る舞います。
豊富な経験とデータに基づいた的確な判断で、部下の成長を支援しながら業務をサポートしてください。

添付の決裁書をレビューする際は、まず作成者をねぎらうために一言褒めてあげてください。
レビュー結果は、見やすいように見出しの先頭に絵文字をつけてください。

【当たり前品質のチェック】
- 体言止めの統一、ページ番号の記載などは抜かりないか。
- 資料はまず冒頭で「何の決裁が欲しいのか」を明確にしてほしい。目的が不明なまま説明が始まると良くない。
- 「良い資料」の条件は、ひとりよがりにならず、「読む人の立場とレベル」を理解して書かれていることである。
- 説明の構成と資料の順番が合っているか確認してほしい (当日、スライドをあちこち移動しなくて済むように練ってほしい)
- 最も重視する要素の1つが「コストの妥当性」
- 常に「ユーザー目線」を優先し、インフラ部門(自分たち)の都合でユーザーに迷惑がかかるような提案は認めない。

【高度な分析】
- 過去の類似案件との比較：同様のシステム導入や改修案件があれば、その成功/失敗要因を参考に助言する
- 業界動向・技術トレンド：最新の技術動向や競合他社の事例を踏まえた妥当性を評価
- リスク分析：技術的リスク、セキュリティリスク、運用リスクを多角的に評価
- グループ社内での整合性：グループ内の他部門や関連会社での類似取り組みや対応事例、連携可能性を検討

【追加のアドバイス】
- 上記以外にも、一般的なレビュー準備の観点で必要と思われる点があれば、追加でアドバイスをしてください。
- 判断根拠を明確に示し、「なぜそう判断したか」を具体的に説明してください。
- 改善提案をする際は、実現可能性を考慮した現実的なアドバイスを心がけてください。
- 案件の実施効果を、定量・定性の両面で例示してください。

【プレゼン時の時間配分】
- 決裁説明時の時間配分のアドバイスも添えてください。10分の会議時間の場合、各セクションの所要時間目安と、質疑応答の時間も加味してください。
- 重要な論点に時間を多く配分し、詳細すぎる技術説明は別資料に回すなど、メリハリのある構成を提案してください。
- これらを表形式で表示してください。

【レビュー後のフォロー】
- 承認条件や宿題事項があれば明確に示してください。
- 次回のレビューに向けた改善ポイントを建設的に伝えてください。

【決裁書内容】
{document_text}"""

def extract_keywords_with_sonnet(bedrock_client, document_text):
    """Claude Sonnet 4を使用して文書から検索キーワードを抽出"""
    try:
        keyword_extraction_prompt = KEYWORD_EXTRACTION_PROMPT_TEMPLATE.format(
            document_text=document_text[:1500]
        )
        
        messages = [
            {
                "role": "user",
                "content": [
                    {
                        "text": keyword_extraction_prompt
                    }
                ]
            }
        ]
        
        response = bedrock_client.converse(
            modelId="us.anthropic.claude-sonnet-4-20250514-v1:0",  # Sonnet 4
            messages=messages,
            inferenceConfig={
                "maxTokens": 4096
            }
        )
        
        # レスポンスから検索キーワードを抽出
        response_text = response['output']['message']['content'][0]['text']
        keywords = []
        
        for line in response_text.split('\n'):
            if 'キーワード' in line and ':' in line:
                keyword = line.split(':', 1)[1].strip()
                if keyword:
                    keywords.append(keyword)
        
        return keywords[:3]  # 最大3個
        
    except Exception as e:
        st.warning(f"キーワード抽出エラー: {e}")
        return ["決裁書", "承認", "ガイドライン"]  # フォールバック

def search_related_information(tavily_client, bedrock_client, document_text, enable_search=True):
    """文書内容に関連する最新情報を検索"""
    if not enable_search or not tavily_client or not bedrock_client:
        return ""
    
    try:
        # Claude Sonnet 4でキーワード抽出
        st.info("検索キーワードを抽出中...")
        extracted_keywords = extract_keywords_with_sonnet(bedrock_client, document_text)
        
        if extracted_keywords:
            st.success(f"✅ 抽出されたキーワード: {', '.join(extracted_keywords)}")
        
        search_results = []
        for keyword in extracted_keywords[:3]:  # 最大3つのキーワードで検索
            try:
                response = tavily_client.search(
                    query=keyword,
                    search_depth="basic",
                    max_results=20,  # 結果数を20件
                    include_answer=True
                )
                
                if response.get('results'):
                    for result in response['results']:
                        search_results.append({
                            'title': result.get('title', ''),
                            'content': result.get('content', ''),
                            'url': result.get('url', ''),
                            'keyword': keyword
                        })
            except Exception as e:
                st.warning(f"検索キーワード '{keyword}' でエラー: {e}")
                continue
        
        # 検索結果をフォーマット（文字数制限付き）
        if search_results:
            formatted_results = "\n\n=== 関連情報（AI抽出キーワード検索結果） ===\n"
            
            for i, result in enumerate(search_results[:5], 1):  # 最大5件に制限
                formatted_results += f"\n{i}. {result['title'][:80]}...\n"  # タイトルも短縮
                formatted_results += f"内容: {result['content'][:150]}...\n"  # 内容をさらに短縮
                formatted_results += f"出典: {result['url']}\n"
                formatted_results += f"検索キーワード: {result['keyword']}\n"
            
            return formatted_results
        else:
            return ""
            
    except Exception as e:
        st.warning(f"関連情報検索エラー: {e}")
        return ""

def create_powerpoint_from_review(review_text, filename="review_result"):
    """レビュー結果からPowerPointプレゼンテーションを作成"""
    try:
        # 新しいプレゼンテーションを作成
        prs = Presentation()
        
        # タイトルスライドを作成
        slide_layout = prs.slide_layouts[0]  # タイトルスライドレイアウト
        slide = prs.slides.add_slide(slide_layout)
        
        # タイトル設定（存在する場合のみ）
        if slide.shapes.title:
            slide.shapes.title.text = "📋 決裁書レビュー結果"
            # タイトルのフォントサイズを設定
            title_text_frame = slide.shapes.title.text_frame
            for paragraph in title_text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(36)
                    run.font.bold = True
        
        # サブタイトル設定（プレースホルダーが存在する場合のみ）
        if len(slide.placeholders) > 1:
            subtitle_placeholder = slide.placeholders[1]
            subtitle_placeholder.text = f"AI部長によるレビュー\n生成日時: {datetime.now().strftime('%Y年%m月%d日 %H:%M')}"
            # サブタイトルのフォントサイズを設定
            subtitle_text_frame = subtitle_placeholder.text_frame
            for paragraph in subtitle_text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(18)
        
        # レビュー内容を解析してセクションに分割
        sections = parse_review_content(review_text)
        
        # 各セクションを別スライドに
        for section_title, content in sections:
            slide_layout = prs.slide_layouts[1]  # タイトルとコンテンツレイアウト
            slide = prs.slides.add_slide(slide_layout)
            
            # タイトル設定
            if slide.shapes.title:
                slide.shapes.title.text = section_title
                # タイトルのフォント設定
                title_text_frame = slide.shapes.title.text_frame
                for paragraph in title_text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(28)
                        run.font.bold = True
            
            # コンテンツ設定（プレースホルダーが存在する場合のみ）
            if len(slide.placeholders) > 1:
                content_placeholder = slide.placeholders[1]
                
                # 長すぎるコンテンツを調整
                if len(content) > 800:
                    content = content[:800] + "\n\n（以下省略）"
                
                content_placeholder.text = content
                
                # フォントサイズを調整
                content_text_frame = content_placeholder.text_frame
                for paragraph in content_text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(16)
                
                # 自動サイズ調整を有効化
                content_text_frame.auto_size = True
        
        # メモリ上でファイルを作成
        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
        
        return ppt_io.getvalue()
        
    except Exception as e:
        st.error(f"PowerPoint生成エラー: {e}")
        return None

def parse_review_content(review_text):
    """レビューテキストを解析してセクションに分割"""
    import re
    from datetime import datetime
    
    sections = []
    
    # 見出しで分割（絵文字付きの見出しを検出）
    lines = review_text.split('\n')
    current_section_title = ""
    current_content = []
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # 見出しの検出（##で始まるもの、または絵文字で始まるもの）
        if (line.startswith('##') or 
            re.match(r'^[🔍📊💡⚠️📋✅❌💰🎯📈📉⏰🔧🚀📝💪👍]+\s*[^\s]', line)):
            
            # 前のセクションを保存
            if current_section_title and current_content:
                sections.append((current_section_title, '\n'.join(current_content)))
            
            # 新しいセクション開始
            current_section_title = line.replace('##', '').strip()
            current_content = []
        else:
            current_content.append(line)
    
    # 最後のセクションを追加
    if current_section_title and current_content:
        sections.append((current_section_title, '\n'.join(current_content)))
    
    # セクションが見つからない場合は全体を一つのスライドに
    if not sections:
        sections = [("レビュー結果", review_text)]
    
    return sections

def sanitize_text_safe_encoding(text):
    """安全なエンコーディング方式でテキストをサニタイズ"""
    if not text:
        return text
    
    import re
    
    try:
        # 方法1: 安全なASCII文字のみ保持（最もシンプル）
        # 英数字、日本語、基本的な句読点のみ許可
        safe_text = re.sub(r'[^\w\s\u3040-\u309F\u30A0-\u30FF\u4E00-\u9FAF\.\,\!\?\:\;\-\(\)\[\]\"\'\/]', ' ', text)
        
        # 複数の空白を単一に
        safe_text = re.sub(r'\s+', ' ', safe_text).strip()
        
        # 長すぎる場合は切り詰め
        if len(safe_text) > 8000:
            safe_text = safe_text[:8000] + "...(省略)"
        
        # ASCII互換性チェック
        try:
            safe_text.encode('ascii', errors='ignore').decode('ascii', errors='ignore')
        except:
            # ASCII化できない場合の代替処理
            safe_text = re.sub(r'[^\x20-\x7E\u3040-\u309F\u30A0-\u30FF\u4E00-\u9FAF]', ' ', safe_text)
            safe_text = re.sub(r'\s+', ' ', safe_text).strip()
        
        return safe_text
        
    except Exception as e:
        # 全ての処理が失敗した場合の最終手段
        st.warning(f"テキスト処理で問題が発生しました: {e}")
        # 最低限の文字のみ保持
        fallback_text = re.sub(r'[^\w\s]', ' ', str(text))
        return re.sub(r'\s+', ' ', fallback_text).strip()[:5000]

def create_review_prompt(document_text, custom_prompt_template, search_results="", additional_message=""):
    """決裁書レビュー用のプロンプトを作成（安全なエンコーディング付き）"""
    # 新しい安全なサニタイズ方式を適用
    document_text = sanitize_text_safe_encoding(document_text)
    
    enhanced_document_text = document_text
    if search_results:
        search_results = sanitize_text_safe_encoding(search_results)
        enhanced_document_text = document_text + search_results
    
    # 追加メッセージがある場合はプロンプトに含める
    if additional_message and additional_message.strip():
        additional_instruction = f"\n\n【追加のレビュー指示】\n{sanitize_text_safe_encoding(additional_message)}\n上記の指示を特に重視してレビューを行ってください。"
        prompt = custom_prompt_template.format(document_text=enhanced_document_text) + additional_instruction
    else:
        prompt = custom_prompt_template.format(document_text=enhanced_document_text)
    
    return prompt

def stream_bedrock_response(bedrock_client, prompt):
    """Bedrock APIを使用してストリーミングレスポンスを生成"""
    try:
        model_id = "us.anthropic.claude-sonnet-4-20250514-v1:0"
        
        messages = [
            {
                "role": "user",
                "content": [
                    {
                        "text": prompt
                    }
                ]
            }
        ]
        
        response = bedrock_client.converse_stream(
            modelId=model_id,
            messages=messages,
            inferenceConfig={
                "maxTokens": 4000
            }
        )
        
        return response
        
    except Exception as e:
        error_msg = str(e)
        if "ServiceUnavailableException" in error_msg:
            st.error("🚫 Bedrock APIが一時的に利用できません。検索機能をオフにするか、より短い文書でお試しください。")
        elif "ThrottlingException" in error_msg:
            st.error("⏱️ APIのリクエスト制限に達しました。しばらく待ってから再試行してください。")
        elif "AccessDeniedException" in error_msg:
            st.error("🔑 AWS認証情報またはモデルアクセス権限を確認してください。")
        else:
            st.error(f"❌ Bedrock API呼び出しエラー: {e}")
        return None

def check_authentication():
    """認証チェック関数"""
    if 'authenticated' not in st.session_state:
        st.session_state.authenticated = False
    
    if not st.session_state.authenticated:
        # 中央寄せのためのレイアウト
        _, col2, _ = st.columns([1, 2, 1])
        
        with col2:
            st.title("🔐 ログイン")
            st.markdown("認証が必要です")
            
            with st.form("login_form"):
                username = st.text_input("ユーザー名")
                password = st.text_input("パスワード", type="password")
                submit_button = st.form_submit_button("ログイン", use_container_width=True)
                
                if submit_button:
                    try:
                        # secrets.tomlから認証情報を取得
                        correct_username = st.secrets["auth"]["username"]
                        correct_password = st.secrets["auth"]["password"]
                        
                        if username == correct_username and password == correct_password:
                            st.session_state.authenticated = True
                            st.success("✅ ログイン成功！")
                            st.rerun()
                        else:
                            st.error("❌ ユーザー名またはパスワードが間違っています")
                    except Exception as e:
                        st.error(f"❌ 認証設定エラー: {e}")
                        st.info("💡 .streamlit/secrets.toml ファイルに認証情報を設定してください")
        
        return False
    
    return True

def main():
    st.set_page_config(
        page_title="AI上司",
        page_icon="👨‍💼",
        layout="wide"
    )
    
    # カスタムCSS
    st.markdown("""
    <style>
    /* 全体のフォント設定 */
    .main {
        padding-top: 2rem;
    }
    
    /* タイトルのスタイリング */
    .main-title {
        font-size: 3rem;
        font-weight: 700;
        color: #1f2937;
        text-align: center;
        margin-bottom: 0.5rem;
    }
    
    .subtitle {
        font-size: 1.2rem;
        color: #6b7280;
        text-align: center;
        margin-bottom: 2rem;
    }
    
    /* サイドバーのスタイリング */
    .css-1d391kg {
        background-color: #f8fafc;
        border-right: 1px solid #e2e8f0;
    }
    
    /* サイドバー全体の余白調整 */
    .stSidebar > div {
        padding-left: 1rem;
        padding-right: 1rem;
    }
    
    /* より具体的なサイドバーのスタイリング */
    [data-testid="stSidebar"] {
        padding-left: 1rem;
        padding-right: 1rem;
    }
    
    [data-testid="stSidebar"] > div {
        padding-left: 0.5rem;
        padding-right: 0.5rem;
    }
    
    /* ボタンのスタイリング */
    .stButton > button {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        border: none;
        border-radius: 8px;
        padding: 0.75rem 1.5rem;
        font-weight: 600;
        transition: all 0.3s ease;
        width: 100%;
    }
    
    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 12px rgba(102, 126, 234, 0.3);
    }
    
    /* プライマリボタン */
    div[data-testid="stButton"] > button[kind="primary"] {
        background: linear-gradient(135deg, #10b981 0%, #059669 100%);
    }
    
    div[data-testid="stButton"] > button[kind="primary"]:hover {
        box-shadow: 0 4px 12px rgba(16, 185, 129, 0.3);
    }
    
    /* アップロードエリアのスタイリング */
    .uploadedFile {
        border: 2px dashed #d1d5db;
        border-radius: 8px;
        padding: 2rem;
        text-align: center;
        background-color: #f9fafb;
        transition: all 0.3s ease;
    }
    
    .uploadedFile:hover {
        border-color: #667eea;
        background-color: #f0f4ff;
    }
    
    /* 成功メッセージのスタイリング */
    .stSuccess > div {
        background-color: #ecfdf5;
        border: 1px solid #a7f3d0;
        border-radius: 8px;
        color: #065f46;
    }
    
    /* エラーメッセージのスタイリング */
    .stError > div {
        background-color: #fef2f2;
        border: 1px solid #fca5a5;
        border-radius: 8px;
        color: #991b1b;
    }
    
    /* 情報メッセージのスタイリング */
    .stInfo > div {
        background-color: #eff6ff;
        border: 1px solid #93c5fd;
        border-radius: 8px;
        color: #1e40af;
    }
    
    /* スピナーのスタイリング */
    .stSpinner > div {
        border-color: #667eea !important;
    }
    </style>
    """, unsafe_allow_html=True)
    
    # 認証チェック
    if not check_authentication():
        return
    
    
    # タイトルエリア
    st.markdown('<h1 class="main-title">👨‍💼 AI上司</h1>', unsafe_allow_html=True)
    st.markdown('<p class="subtitle">あなたの上司に代わって、決裁資料のレビューをします</p>', unsafe_allow_html=True)
    
    # サイドバーでレビュープロンプト設定
    with st.sidebar:
        st.header("⚙️ レビュー設定")
        
        # デフォルトプロンプトテンプレート（変数から取得）
        default_prompt = DEFAULT_REVIEW_PROMPT_TEMPLATE
        
        custom_prompt = st.text_area(
            "レビュープロンプト",
            value=default_prompt,
            height=500,
            help="プロンプト内で {document_text} を使用すると、PDFの内容が挿入されます"
        )
        
        if st.button("プロンプトをリセット"):
            st.session_state.custom_prompt = DEFAULT_REVIEW_PROMPT_TEMPLATE
            st.rerun()
            
        # セッション状態にプロンプトを保存
        st.session_state.custom_prompt = custom_prompt
        
        st.divider()
        
        # 検索オプション
        st.subheader("🔍 検索設定")
        enable_search = st.checkbox(
            "関連情報をWeb検索してレビュー品質を向上",
            value=True,
            help="Tavily APIを使用して、決裁書に関連する最新情報を検索し、レビューの参考にします"
        )
    
    # メインエリア    
    uploaded_file = st.file_uploader(
        "決裁書をアップロードしてください",
        type=['pdf', 'pptx'],
        help="PDFファイル(.pdf)またはPowerPointファイル(.pptx)に対応しています"
    )
    
    if uploaded_file is not None:
        st.success(f"✅ ファイルアップロード完了")
        
        # 追加のメッセージ入力フォーム
        st.markdown("### 📝 追加のレビュー指示（任意）")
        st.markdown("*今回のレビューで特に注意したい点があれば、下記に入力してください*")
        additional_message = st.text_area(
            "追加の指示内容",
            placeholder="例：今回のレビューは定性効果のみ訴求すればOKで、定量効果は不要です",
            help="空欄のままでも問題ありません。必要に応じて追加の指示を記入してください。",
            height=100,
            label_visibility="collapsed"
        )
        
        # ファイル形式を判定してテキスト抽出
        file_extension = uploaded_file.name.lower().split('.')[-1]
        
        if file_extension == 'pdf':
            with st.spinner("PDF内容を読み込み中..."):
                document_text = extract_text_from_pdf(uploaded_file)
        elif file_extension == 'pptx':
            with st.spinner("PowerPoint内容を読み込み中..."):
                document_text = extract_text_from_pptx(uploaded_file)
        else:
            st.error(f"サポートされていないファイル形式です: {file_extension}")
            document_text = None
        
        if document_text:            
            # レビュー実行ボタン
            if st.button("🔍 AIレビューを開始", type="primary"):
                bedrock_client = init_bedrock_client()
                
                if bedrock_client:
                    search_results = ""
                    
                    # 関連情報検索（有効な場合）
                    if enable_search:
                        with st.spinner("関連情報を検索中..."):
                            tavily_client = init_tavily_client()
                            if tavily_client:
                                search_results = search_related_information(tavily_client, bedrock_client, document_text, enable_search)
                                if search_results:
                                    st.success("✅ 関連情報の検索完了")
                                    with st.expander("🔎 検索された関連情報"):
                                        st.markdown(search_results)
                                else:
                                    st.info("ℹ️ 追加の関連情報は見つかりませんでした")
                    
                    # プロンプト作成
                    prompt = create_review_prompt(document_text, st.session_state.get('custom_prompt', ''), search_results, additional_message)
                    
                    # ストリーミングレスポンス表示
                    with st.spinner("AIレビューを実行中..."):
                        response_stream = stream_bedrock_response(bedrock_client, prompt)
                        
                        if response_stream:
                            # ストリーミング結果を表示するコンテナ
                            response_container = st.empty()
                            full_response = ""
                            
                            try:
                                for event in response_stream['stream']:
                                    if 'contentBlockDelta' in event:
                                        delta = event['contentBlockDelta']['delta']
                                        if 'text' in delta:
                                            full_response += delta['text']
                                            response_container.markdown(full_response)
                                
                                # 最終結果の保存オプション
                                st.success("✅ レビュー完了")
                                
                                # PowerPointダウンロードボタン
                                ppt_data = create_powerpoint_from_review(full_response)
                                if ppt_data:
                                    st.download_button(
                                        label="📊 レビュー結果をPowerPointでダウンロード",
                                        data=ppt_data,
                                        file_name=f"review_{uploaded_file.name.replace('.pdf', '').replace('.pptx', '')}.pptx",
                                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                        type="primary"
                                    )
                                
                            except Exception as e:
                                st.error(f"ストリーミング処理エラー: {e}")

if __name__ == "__main__":
    main()